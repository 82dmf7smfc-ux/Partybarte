"""End-to-end test.

Runs the whole pipeline through main.run and checks that both files are written
and that they open. Opening the files proves they are valid, not just present.
"""

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from alarm_pareto import main
from tests import data_paths as dp


class _Args:
    """A tiny stand-in for the parsed command line arguments."""

    def __init__(self, output_dir, start_time=None, end_time=None,
                 downtime_method="attributed"):
        self.input = str(dp.SAMPLE_CSV)
        self.vendor = "amat"
        self.config = str(dp.CONFIG_PATH)
        self.window_days = 30
        self.start_time = start_time
        self.end_time = end_time
        self.top_n = 15
        self.downtime_method = downtime_method
        self.output_dir = str(output_dir)


def test_pipeline_writes_and_opens_both_files(tmp_path):
    out = main.run(_Args(tmp_path))

    assert out["xlsx"].exists()
    assert out["pptx"].exists()

    # The workbook must open and have the expected sheets.
    wb = load_workbook(out["xlsx"])
    for sheet in ["Window_Data", "By_Fault_Code", "By_Description", "By_Equipment"]:
        assert sheet in wb.sheetnames

    # Each summary sheet must carry at least one native chart object.
    assert len(wb["By_Fault_Code"]._charts) >= 1

    # The deck must open and have five slides: title, three levels, summary.
    prs = Presentation(out["pptx"])
    assert len(prs.slides) == 5


def test_headline_numbers_are_stable(tmp_path):
    expected = dp.load_expected()["grand"]
    out = main.run(_Args(tmp_path))
    grand = out["result"]["grand"]
    assert grand["total_faults"] == expected["total_faults"]
    assert grand["attributed_downtime_s"] == expected["attributed_downtime_s"]
    assert grand["wallclock_downtime_s"] == expected["wallclock_downtime_s"]


def test_time_of_day_filter_narrows_the_report(tmp_path):
    """A day shift and a night shift must add back up to the whole window."""
    whole = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    day = main.run(_Args(tmp_path / "day", "06:00", "18:00"))["result"]["grand"]
    night = main.run(_Args(tmp_path / "night", "18:00", "06:00"))["result"]["grand"]

    assert day["total_faults"] + night["total_faults"] == whole["total_faults"]
    assert day["total_faults"] < whole["total_faults"]


def test_time_of_day_label_reaches_the_headline_numbers(tmp_path):
    grand = main.run(_Args(tmp_path / "night", "22:00", "06:00"))["result"]["grand"]
    assert grand["time_of_day_label"] == "22:00 to 06:00 (crosses midnight)"

    grand = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    assert grand["time_of_day_label"] == "All hours"


def test_time_of_day_does_not_move_the_window(tmp_path):
    """The window is trimmed first, so the shift never shifts its edges."""
    whole = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    day = main.run(_Args(tmp_path / "day", "06:00", "18:00"))["result"]["grand"]
    assert day["window_start"] == whole["window_start"]
    assert day["window_end"] == whole["window_end"]


def test_one_time_bound_alone_is_an_error(tmp_path):
    with pytest.raises(ValueError):
        main.run(_Args(tmp_path / "half", "06:00", None))


def test_a_range_with_no_alarms_is_an_error(tmp_path):
    with pytest.raises(ValueError):
        main.run(_Args(tmp_path / "quiet", "01:00", "02:00"))


def test_filtered_run_still_writes_both_files(tmp_path):
    out = main.run(_Args(tmp_path / "night", "18:00", "06:00"))
    assert out["xlsx"].exists()
    assert out["pptx"].exists()
    assert len(Presentation(out["pptx"]).slides) == 5


# --- the third downtime number, end to end --------------------------------

def test_all_three_downtime_numbers_are_reported(tmp_path):
    grand = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    for key in ["attributed_downtime_hours", "wallclock_downtime_hours",
                "in_range_downtime_hours", "range_hours", "in_range_downtime_pct"]:
        assert key in grand
        assert grand[key] >= 0


def test_attributed_is_at_least_wall_clock(tmp_path):
    """Crediting every fault its full duration cannot come out under the merge."""
    grand = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    assert grand["attributed_downtime_hours"] >= grand["wallclock_downtime_hours"]


def test_in_range_never_exceeds_the_time_covered(tmp_path):
    """The tool cannot be down for longer than the report is long."""
    for name, args in [("all", ()), ("night", ("18:00", "06:00")), ("day", ("06:00", "18:00"))]:
        grand = main.run(_Args(tmp_path / name, *args))["result"]["grand"]
        assert grand["in_range_downtime_hours"] <= grand["range_hours"] + 1e-9
        assert 0 <= grand["in_range_downtime_pct"] <= 100


def test_day_and_night_in_range_downtime_add_up_to_the_whole_window(tmp_path):
    """The property that makes this number worth having.

    Split the window into two shifts and the in-range downtime splits with it.
    Neither of the other two numbers does this, because they follow the fault
    onset rather than the clock.
    """
    whole = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    day = main.run(_Args(tmp_path / "day", "06:00", "18:00"))["result"]["grand"]
    night = main.run(_Args(tmp_path / "night", "18:00", "06:00"))["result"]["grand"]

    total = day["in_range_downtime_s"] + night["in_range_downtime_s"]
    assert total == pytest.approx(whole["in_range_downtime_s"])
    # And the hours covered split the same way.
    assert day["range_hours"] + night["range_hours"] == pytest.approx(whole["range_hours"])


def test_attributed_downtime_also_splits_by_shift(tmp_path):
    """Attributed splits too, but for a different reason.

    Attributed downtime is conserved per row, and the shift filter puts every
    row in exactly one shift, so it adds up. What it is not is bounded by the
    clock: a shift can be credited more downtime than the shift is long. Only
    the in-range number is bounded, which is checked separately above.
    """
    whole = main.run(_Args(tmp_path / "all"))["result"]["grand"]
    day = main.run(_Args(tmp_path / "day", "06:00", "18:00"))["result"]["grand"]
    night = main.run(_Args(tmp_path / "night", "18:00", "06:00"))["result"]["grand"]

    assert (day["attributed_downtime_s"] + night["attributed_downtime_s"]
            == pytest.approx(whole["attributed_downtime_s"]))


def test_the_three_numbers_on_the_sample_are_exactly_these(tmp_path):
    """Golden values, worked out by hand from the fourteen rows in the window.

    An inequality can pass for the wrong reason. These are the actual numbers,
    so a change in any of the three shows up here as a number, not a hunch.

    The window runs 2026-01-22 20:00 to 2026-02-21 20:00, 720 hours.

    Day shift keeps eleven rows totalling 15.5 hours of duration. Two of them
    overlap (10:00-14:00 and 12:00-15:00 on the 10th), so wall clock drops to
    13.5. Nothing on day shift runs outside the covered hours, so in-range is
    13.5 as well.

    Night shift keeps three rows totalling 3.0 hours with no overlaps, so wall
    clock is also 3.0. But the last of them starts at 20:00 on the 21st, which
    is the exact moment the window closes, so all thirty minutes of it fall
    outside. In-range is therefore 2.5.
    """
    expected = {
        #                    attributed  wallclock  in_range  covered
        "all":   (18.5, 16.5, 16.0, 720.0),
        "day":   (15.5, 13.5, 13.5, 360.0),
        "night": (3.0,  3.0,  2.5,  360.0),
    }
    runs = {
        "all": _Args(tmp_path / "all"),
        "day": _Args(tmp_path / "day", "06:00", "18:00"),
        "night": _Args(tmp_path / "night", "18:00", "06:00"),
    }
    for name, args in runs.items():
        grand = main.run(args)["result"]["grand"]
        attributed, wallclock, in_range, covered = expected[name]
        assert grand["attributed_downtime_hours"] == pytest.approx(attributed), name
        assert grand["wallclock_downtime_hours"] == pytest.approx(wallclock), name
        assert grand["in_range_downtime_hours"] == pytest.approx(in_range), name
        assert grand["range_hours"] == pytest.approx(covered), name


def test_an_alarm_running_past_the_window_end_is_clipped(tmp_path):
    """The reason night in-range is 2.5 and not 3.0, stated on its own.

    The sample's last alarm starts at 20:00 on the 21st and runs half an hour.
    The window closes at 20:00 on the 21st. Wall clock credits the full half
    hour because the alarm started inside the window. In-range credits none of
    it, because none of it happened inside the window.
    """
    night = main.run(_Args(tmp_path / "night", "18:00", "06:00"))["result"]["grand"]
    assert night["wallclock_downtime_hours"] == pytest.approx(3.0)
    assert night["in_range_downtime_hours"] == pytest.approx(2.5)


def test_ranking_by_in_range_works_end_to_end(tmp_path):
    out = main.run(_Args(tmp_path / "night", "18:00", "06:00", "in_range"))
    assert out["xlsx"].exists()
    assert out["pptx"].exists()
    assert out["result"]["downtime_method"] == "in_range"

    by_downtime = out["result"]["levels"]["fault_code"]["by_downtime"]
    values = by_downtime["in_range_s"].tolist()
    assert values == sorted(values, reverse=True)


def test_every_level_table_carries_all_three_measures(tmp_path):
    result = main.run(_Args(tmp_path / "night", "18:00", "06:00"))["result"]
    for level in ["fault_code", "description", "equipment"]:
        table = result["levels"][level]["by_count"]
        for column in ["attributed_hours", "wallclock_hours", "in_range_hours"]:
            assert column in table.columns


def test_a_bad_downtime_method_is_rejected(tmp_path):
    with pytest.raises(ValueError):
        main.run(_Args(tmp_path / "bad", None, None, "nonsense"))


def test_reporting_range_blocks_match_the_shift(tmp_path):
    """No shift means one block. A night shift means one block per night."""
    whole = main.run(_Args(tmp_path / "all"))["result"]
    assert whole["grand"]["range_blocks"] == 1
    assert len(whole["reporting_ranges"]) == 1

    night = main.run(_Args(tmp_path / "night", "18:00", "06:00"))["result"]
    assert night["grand"]["range_blocks"] > 1
    # Every night is twelve hours except the two clipped by the window ends.
    assert night["grand"]["range_hours"] < whole["grand"]["range_hours"]
