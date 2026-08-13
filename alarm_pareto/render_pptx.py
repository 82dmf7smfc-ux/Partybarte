"""Step 6: write the PowerPoint deck.

Excel needs live charts you can click into. A slide deck does not. So here we
draw each Pareto with matplotlib and drop the picture onto a slide. This keeps
the deck simple and fully offline.

The deck has one Pareto slide per grouping level, plus a summary slide with the
headline numbers.
"""

import io

import matplotlib

# Use a non-interactive backend. This draws to memory, never to a screen. It
# also means the code runs on a machine with no display attached.
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402  (import after backend is set)
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from . import aggregate as agg  # noqa: E402

SLIDE_WIDTH = Inches(13.333)   # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

NICE_LEVEL = {"fault_code": "Fault Code", "description": "Description", "equipment": "Equipment"}


def write_deck(result, output_path):
    """Write the full deck to output_path."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    method = result["downtime_method"]
    method_label = _method_label(method)

    _add_title_slide(prs, result, method_label)
    for level in agg.GROUPING_LEVELS:
        _add_pareto_slide(prs, result, level, method, method_label)
    _add_summary_slide(prs, result, method_label)

    prs.save(output_path)
    return output_path


def _method_label(method):
    if method == agg.METHOD_ATTRIBUTED:
        return "Attributed downtime (each fault gets its full duration)"
    return "True wall-clock downtime (overlaps merged)"


def _blank_slide(prs):
    """Return a slide with no placeholders so we can position freely."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
    return box


def _add_title_slide(prs, result, method_label):
    slide = _blank_slide(prs)
    grand = result["grand"]
    _add_textbox(slide, Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
                 "Alarm Pareto Report", size=40, bold=True)
    subtitle = "Window: %s to %s  (%s days)" % (
        _fmt_ts(grand["window_start"]), _fmt_ts(grand["window_end"]), grand["window_days"])
    _add_textbox(slide, Inches(0.7), Inches(3.6), Inches(12), Inches(0.7), subtitle, size=20)
    _add_textbox(slide, Inches(0.7), Inches(4.3), Inches(12), Inches(0.7),
                 "Downtime ranking method: %s" % method_label, size=16)


def _add_pareto_slide(prs, result, level, method, method_label):
    """One slide per level. Shows the count Pareto and the downtime Pareto."""
    slide = _blank_slide(prs)
    tables = result["levels"][level]

    _add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6),
                 "Pareto by %s" % NICE_LEVEL[level], size=26, bold=True)
    _add_textbox(slide, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4),
                 "Downtime method: %s" % method_label, size=12)

    # Left chart: by count. Right chart: by downtime.
    count_png = _pareto_png(
        tables["by_count"], level,
        value_col="count", cum_col="cum_count_pct",
        title="By occurrence count", value_axis="Count",
    )
    metric_hours = "attributed_hours" if method == agg.METHOD_ATTRIBUTED else "wallclock_hours"
    downtime_png = _pareto_png(
        tables["by_downtime"], level,
        value_col=metric_hours, cum_col="cum_downtime_pct",
        title="By downtime (%s)" % method, value_axis="Downtime (hours)",
    )

    if count_png is not None:
        slide.shapes.add_picture(count_png, Inches(0.3), Inches(1.4), width=Inches(6.3))
    if downtime_png is not None:
        slide.shapes.add_picture(downtime_png, Inches(6.8), Inches(1.4), width=Inches(6.3))


def _add_summary_slide(prs, result, method_label):
    slide = _blank_slide(prs)
    grand = result["grand"]
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Summary", size=30, bold=True)

    lines = []
    lines.append("Date range covered: %s to %s (%s days)" % (
        _fmt_ts(grand["window_start"]), _fmt_ts(grand["window_end"]), grand["window_days"]))
    lines.append("Total faults in window: %d" % grand["total_faults"])
    lines.append("")
    lines.append("Total downtime, both methods (never mix these):")
    lines.append("  Attributed: %.2f hours" % grand["attributed_downtime_hours"])
    lines.append("  True wall clock: %.2f hours" % grand["wallclock_downtime_hours"])
    lines.append("")
    lines.append("Top three offenders by %s, at fault-code level:" % result["downtime_method"])
    if result["top_offenders"]:
        for i, off in enumerate(result["top_offenders"], start=1):
            lines.append("  %d. %s  |  %d faults  |  %.2f h attributed  |  %.2f h wall clock" % (
                i, off["name"], off["count"], off["attributed_hours"], off["wallclock_hours"]))
    else:
        lines.append("  No faults in the window.")

    box = _add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
                       "\n".join(lines), size=16)
    # Give the body a little more room per line.
    for para in box.text_frame.paragraphs:
        para.space_after = Pt(4)


def _pareto_png(table, level, value_col, cum_col, title, value_axis):
    """Draw a standard Pareto chart and return it as a PNG in memory."""
    if table is None or table.empty:
        return None

    labels = [str(x) for x in table[level].tolist()]
    values = table[value_col].astype(float).tolist()
    cum = table[cum_col].astype(float).tolist()

    fig, ax1 = plt.subplots(figsize=(6.3, 4.6), dpi=150)

    positions = range(len(labels))
    ax1.bar(positions, values, color="#305496")
    ax1.set_ylabel(value_axis)
    ax1.set_title(title)
    ax1.set_xticks(list(positions))
    # Long labels are rotated so they do not overlap.
    ax1.set_xticklabels(labels, rotation=45, ha="right", fontsize=7)

    # The cumulative line lives on a second y axis on the right.
    ax2 = ax1.twinx()
    ax2.plot(positions, cum, color="#C00000", marker="o", linewidth=2)
    ax2.set_ylabel("Cumulative %")
    ax2.set_ylim(0, 105)

    # The 80 percent reference line. Dashed so it reads as a guide, not data.
    ax2.axhline(80, color="gray", linestyle="--", linewidth=1)

    fig.tight_layout()

    buffer = io.BytesIO()
    fig.savefig(buffer, format="png")
    plt.close(fig)  # free the figure so memory does not build up
    buffer.seek(0)
    return buffer


def _fmt_ts(ts):
    import pandas as pd
    if ts is None or (isinstance(ts, pd.Timestamp) and pd.isna(ts)):
        return "n/a"
    return pd.Timestamp(ts).strftime("%Y-%m-%d %H:%M:%S")
